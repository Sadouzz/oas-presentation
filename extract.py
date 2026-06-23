from pptx import Presentation
import sys

def extract_text(filepath):
    try:
        prs = Presentation(filepath)
        for i, slide in enumerate(prs.slides):
            print(f"--- Slide {i+1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    print(shape.text)
            print()
    except Exception as e:
        print(f"Error: {e}")

if __name__ == "__main__":
    extract_text("Conception et Développement OAS Web et Billing.pptx")
